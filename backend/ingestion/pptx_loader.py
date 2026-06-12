import logging
from pptx import Presentation
from .base_loader import DocumentLoader

logger = logging.getLogger(__name__)

class PptxLoader(DocumentLoader):
    """Loader for PowerPoint (PPTX) documents."""

    def extract_text(self, file_path: str) -> str:
        """
        Extract text from a PPTX file using python-pptx.
        
        Args:
            file_path (str): Path to the PPTX file.
            
        Returns:
            str: The extracted text.
        """
        logger.info(f"Extracting text from PPTX: {file_path}")
        try:
            prs = Presentation(file_path)
            text_chunks = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_chunks.append(shape.text)
                        
            return "\n".join(text_chunks)
        except Exception as e:
            logger.error(f"Error extracting text from PPTX {file_path}: {e}")
            raise
